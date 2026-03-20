"""Generate the StepBy hackathon 'Logic' slide as a polished PPTX."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x26, 0x22, 0x62)
HEADER_BLUE = RGBColor(0x1B, 0x75, 0xBC)
CYAN = RGBColor(0x00, 0xAE, 0xEF)
GITHUB_BG = RGBColor(0x2D, 0x33, 0x3B)
GEMINI_BG = RGBColor(0x6C, 0x5C, 0xE7)
UI_BG = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFB)
CARD_BG = RGBColor(0xF8, 0xFB, 0xFF)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT_LINE = RGBColor(0xD9, 0xE1, 0xEC)


ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "青と白　シンプル　ポートフォリオ　プレゼンテーション.pptx"
OUTPUT_DIR = ROOT / "outputs" / "StepByハッカソン"
OUTPUT_PATH = OUTPUT_DIR / "ロジック_スライド.pptx"


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=18,
    bold=False,
    color=NAVY,
    alignment=PP_ALIGN.LEFT,
    font_name="Yu Gothic UI",
    vertical_anchor=MSO_ANCHOR.TOP,
    margin_left=0.0,
    margin_right=0.0,
    margin_top=0.0,
    margin_bottom=0.0,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(margin_right)
    tf.margin_top = Inches(margin_top)
    tf.margin_bottom = Inches(margin_bottom)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    paragraphs,
    *,
    color=WHITE,
    font_name="Yu Gothic UI",
    margin=0.16,
    vertical_anchor=MSO_ANCHOR.MIDDLE,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)

    for idx, item in enumerate(paragraphs):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.alignment = item.get("align", PP_ALIGN.CENTER)
        paragraph.space_after = Pt(item.get("space_after", 4))
        paragraph.space_before = Pt(item.get("space_before", 0))
        run = paragraph.add_run()
        run.text = item["text"]
        run.font.size = Pt(item.get("size", 16))
        run.font.bold = item.get("bold", False)
        run.font.color.rgb = item.get("color", color)
        run.font.name = item.get("font_name", font_name)
    return box


def add_round_rect(slide, left, top, width, height, fill, *, radius=0.08, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
      shape.line.color.rgb = line
      shape.line.width = Pt(1.2)
    else:
      shape.line.fill.background()
    shape.adjustments[0] = radius
    return shape


def add_circle(slide, left, top, size, fill):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill
    circle.line.fill.background()
    return circle


def add_vertical_divider(slide, x, top, height):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, top, x, top + height)
    line.line.color.rgb = LIGHT_LINE
    line.line.width = Pt(1.3)
    return line


def add_arrow(slide, left, top, width, height):
    arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = CYAN
    arrow.line.fill.background()
    return arrow


def add_monitor_icon(slide, left, top, size):
    screen = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size * 0.62
    )
    screen.fill.solid()
    screen.fill.fore_color.rgb = WHITE
    screen.line.fill.background()
    screen.adjustments[0] = 0.08

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + size * 0.18,
        top + size * 0.67,
        size * 0.64,
        size * 0.08,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()

    stand = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        left + size * 0.41,
        top + size * 0.57,
        size * 0.18,
        size * 0.18,
    )
    stand.fill.solid()
    stand.fill.fore_color.rgb = WHITE
    stand.line.fill.background()
    stand.rotation = 180


def build_slide():
    src = Presentation(str(TEMPLATE))
    prs = Presentation()
    prs.slide_width = src.slide_width
    prs.slide_height = src.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Background
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE

    # ------------------------------------------------------------------
    # Header
    # ------------------------------------------------------------------
    header_h = Inches(1.65)
    left_block_w = Inches(3.55)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, header_h)
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    left_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(-Inches(0.18)), 0, left_block_w, header_h
    )
    left_block.fill.solid()
    left_block.fill.fore_color.rgb = HEADER_BLUE
    left_block.line.fill.background()

    add_text_box(
        slide,
        Inches(0.48),
        Inches(0.10),
        Inches(1.7),
        Inches(1.3),
        "03",
        font_size=54,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text_box(
        slide,
        Inches(2.62),
        Inches(0.18),
        Inches(7.6),
        Inches(1.2),
        "ロジック",
        font_size=58,
        bold=True,
        color=WHITE,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    # ------------------------------------------------------------------
    # Subtitle
    # ------------------------------------------------------------------
    subtitle_top = Inches(1.95)
    add_text_box(
        slide,
        Inches(1.3),
        subtitle_top,
        Inches(17.4),
        Inches(0.55),
        "GitHub API  →  Gemini AI（JSON）  →  初心者向けUI",
        font_size=24,
        bold=True,
        color=NAVY,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    # ------------------------------------------------------------------
    # Flow area
    # ------------------------------------------------------------------
    flow_top = Inches(2.65)
    flow_h = Inches(3.45)
    col_w = Inches(5.4)
    gap = Inches(0.55)
    start_x = Inches(1.2)
    circle_size = Inches(1.28)
    circle_y = flow_top + Inches(0.08)
    title_y = flow_top + Inches(1.52)
    card_y = flow_top + Inches(2.02)
    card_h = Inches(1.18)

    columns = [
        {
            "x": start_x,
            "circle_fill": CYAN,
            "card_fill": GITHUB_BG,
            "title": "GitHub API / Octokit",
            "body": [
                "リポジトリ情報を取得",
                "・ファイル構造 / README",
                "・Issue（title, body）",
                "・PR diff",
            ],
            "title_color": NAVY,
            "body_color": WHITE,
            "icon": "github",
        },
        {
            "x": start_x + col_w + gap,
            "circle_fill": CYAN,
            "card_fill": GEMINI_BG,
            "title": "Gemini 2.5 Flash",
            "body": [
                "初心者向けに構造化",
                "→ JSON出力",
            ],
            "title_color": NAVY,
            "body_color": WHITE,
            "icon": "sparkle",
        },
        {
            "x": start_x + (col_w + gap) * 2,
            "circle_fill": CYAN,
            "card_fill": UI_BG,
            "title": "初心者向けUI",
            "body": [
                "わかりやすく表示",
                "・極小ステップ",
                "・公式ドキュメントリンク",
                "・コマンドコピペ",
            ],
            "title_color": NAVY,
            "body_color": WHITE,
            "icon": "monitor",
        },
    ]

    for idx, col in enumerate(columns):
        center_x = col["x"] + (col_w - circle_size) / 2
        add_circle(slide, center_x, circle_y, circle_size, col["circle_fill"])

        if col["icon"] == "github":
            add_text_box(
                slide,
                center_x,
                circle_y + Inches(0.05),
                circle_size,
                circle_size - Inches(0.05),
                "GH",
                font_size=28,
                bold=True,
                color=WHITE,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
            )
        elif col["icon"] == "sparkle":
            add_text_box(
                slide,
                center_x,
                circle_y + Inches(0.02),
                circle_size,
                circle_size,
                "✦",
                font_size=34,
                bold=True,
                color=WHITE,
                alignment=PP_ALIGN.CENTER,
                vertical_anchor=MSO_ANCHOR.MIDDLE,
                font_name="Yu Gothic UI Symbol",
            )
        else:
            add_monitor_icon(slide, center_x + Inches(0.18), circle_y + Inches(0.24), Inches(0.92))

        if idx < 2:
            add_vertical_divider(
                slide,
                col["x"] + col_w + gap / 2,
                flow_top + Inches(0.10),
                flow_h - Inches(0.15),
            )

        add_text_box(
            slide,
            col["x"] + Inches(0.05),
            title_y,
            col_w - Inches(0.1),
            Inches(0.40),
            col["title"],
            font_size=23,
            bold=True,
            color=col["title_color"],
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_round_rect(
            slide,
            col["x"] + Inches(0.18),
            card_y,
            col_w - Inches(0.36),
            card_h,
            col["card_fill"],
            radius=0.08,
        )

        paragraphs = [{"text": col["body"][0], "size": 18, "bold": True}]
        for line in col["body"][1:]:
            paragraphs.append({"text": line, "size": 14, "space_after": 2})

        add_paragraphs(
            slide,
            col["x"] + Inches(0.32),
            card_y + Inches(0.06),
            col_w - Inches(0.64),
            card_h - Inches(0.12),
            paragraphs,
            color=col["body_color"],
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

    add_arrow(slide, start_x + col_w + Inches(0.08), flow_top + Inches(1.42), Inches(0.34), Inches(0.24))
    add_arrow(
        slide,
        start_x + (col_w + gap) * 2 - Inches(0.47),
        flow_top + Inches(1.42),
        Inches(0.34),
        Inches(0.24),
    )

    # ------------------------------------------------------------------
    # Pipeline cards
    # ------------------------------------------------------------------
    list_top = Inches(6.45)
    list_left = Inches(1.35)
    list_w = Inches(17.3)
    card_gap = Inches(0.28)
    pipeline_card_w = (list_w - card_gap) / 2
    pipeline_card_h = Inches(0.86)

    pipeline_items = [
        ("①", "リポジトリ分析", "学習ロードマップ + 全Issue分解"),
        ("②", "Issue", "極小タスク分解"),
        ("③", "プロジェクト", "オンボーディングガイド自動生成"),
        ("④", "PR diff", "AIコードレビュー"),
    ]

    for i, (num, left_text, right_text) in enumerate(pipeline_items):
        row = i // 2
        col = i % 2
        x = list_left + col * (pipeline_card_w + card_gap)
        y = list_top + row * (pipeline_card_h + Inches(0.22))

        add_round_rect(slide, x, y, pipeline_card_w, pipeline_card_h, CARD_BG, radius=0.05, line=LIGHT_LINE)
        add_circle(slide, x + Inches(0.18), y + Inches(0.16), Inches(0.52), CYAN)
        add_text_box(
            slide,
            x + Inches(0.18),
            y + Inches(0.14),
            Inches(0.52),
            Inches(0.52),
            num,
            font_size=16,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

        add_text_box(
            slide,
            x + Inches(0.82),
            y + Inches(0.11),
            pipeline_card_w - Inches(1.02),
            Inches(0.24),
            f"{left_text}  →",
            font_size=15,
            bold=True,
            color=NAVY,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text_box(
            slide,
            x + Inches(0.82),
            y + Inches(0.37),
            pipeline_card_w - Inches(1.02),
            Inches(0.26),
            right_text,
            font_size=13,
            bold=False,
            color=MUTED,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

    # ------------------------------------------------------------------
    # Bottom note
    # ------------------------------------------------------------------
    note_top = Inches(8.72)
    add_round_rect(
        slide,
        Inches(1.32),
        note_top,
        Inches(17.36),
        Inches(0.64),
        LIGHT_BG,
        radius=0.04,
        line=LIGHT_LINE,
    )
    add_text_box(
        slide,
        Inches(1.58),
        note_top + Inches(0.06),
        Inches(16.84),
        Inches(0.52),
        "全パイプライン共通:  Gemini 2.5 Flash  /  JSON構造化出力  /  Firestore永続化",
        font_size=15,
        bold=True,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    return prs


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_slide()
    prs.save(str(OUTPUT_PATH))
    print(f"Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
